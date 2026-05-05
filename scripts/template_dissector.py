#!/usr/bin/env python3
"""Template Dissector v2 — 模板解剖器

解剖PPTX模板，逐页提取6类元素，产出JSON + 合入spec_lock/design_spec。
v2: 增强主题色(theme color)和主题字体(theme font)的深度提取。

Usage:
    python template_dissector.py <pptx_path> <project_path>
"""

import win_compat  # noqa: F401 -- Windows compat (path+encoding fix)

import sys
import os
import json
from collections import Counter
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn, nsmap
    import lxml.etree as etree
except ImportError:
    sys.stderr.write("ERROR: python-pptx or lxml not installed. Run: pip install python-pptx lxml\n")
    sys.exit(1)


def emu_to_px(emu):
    """Convert EMU to pixels (1 inch = 914400 EMU, 1 inch = 96 px)"""
    if emu is None:
        return None
    return round(emu / 914400 * 96, 1)


def emu_to_pt(emu):
    """Convert EMU to points (1 inch = 914400 EMU, 1 inch = 72 pt)"""
    if emu is None:
        return None
    return round(emu / 914400 * 72, 1)


def rgb_to_hex(rgb_color):
    """Convert RGBColor to hex string"""
    if rgb_color is None:
        return None
    try:
        return f"#{rgb_color}"
    except Exception:
        return None


# ============================================================
# Theme extraction helpers
# ============================================================

def _parse_theme_xml(theme_part):
    """Parse theme XML from a Part object (uses .blob since .element may not exist)."""
    try:
        # Try .element first (XmlPart subclasses)
        if hasattr(theme_part, 'element'):
            return theme_part.element
    except Exception:
        pass
    
    # Fallback: parse from blob
    try:
        blob = theme_part.blob
        return etree.fromstring(blob)
    except Exception:
        return None


NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}


def extract_theme_colors(prs):
    """Extract the color scheme from the presentation's theme.
    
    Returns a dict mapping scheme color names to hex values,
    e.g. {"dk1": "#000000", "lt1": "#FFFFFF", "dk2": "#44546A", ...}
    """
    theme_colors = {}
    theme_name = ""
    
    try:
        # Navigate to theme via slide master
        for slide_master in prs.slide_masters:
            for rel in slide_master.part.rels.values():
                if "theme" in rel.reltype:
                    theme_part = rel.target_part
                    theme_elem = _parse_theme_xml(theme_part)
                    
                    if theme_elem is None:
                        continue
                    
                    # Find the clrScheme element (it's nested under themeElements, use .//)
                    clr_scheme = theme_elem.find('.//a:clrScheme', NS)
                    if clr_scheme is None:
                        # Fallback: try findall
                        results = theme_elem.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}clrScheme')
                        if results:
                            clr_scheme = results[0]
                    if clr_scheme is None:
                        continue
                    
                    theme_name = clr_scheme.get('name', '')
                    
                    for child in clr_scheme:
                        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                        # Each child has a sub-element like a:sysClr or a:srgbClr
                        for color_elem in child:
                            color_tag = color_elem.tag.split('}')[-1] if '}' in color_elem.tag else color_elem.tag
                            if color_tag == 'srgbClr':
                                val = color_elem.get('val')
                                if val:
                                    theme_colors[tag] = f"#{val}"
                            elif color_tag == 'sysClr':
                                # For sysClr, prefer lastClr (the cached RGB), fallback to val
                                val = color_elem.get('lastClr')
                                if not val:
                                    # Map common sysClr val to actual colors
                                    sys_map = {
                                        'windowText': '000000',  # Dark text
                                        'window': 'FFFFFF',      # Light background
                                    }
                                    raw = color_elem.get('val', '')
                                    val = sys_map.get(raw, raw)
                                if val and not val.startswith('#'):
                                    theme_colors[tag] = f"#{val}"
                    
                    break  # Only process first theme
            break  # Only process first slide master
    except Exception as e:
        sys.stderr.write(f"  [WARN] Theme color extraction failed: {e}\n")
    
    # Store theme name as metadata
    if theme_name:
        theme_colors["_theme_name"] = theme_name
    
    return theme_colors


def extract_theme_fonts(prs):
    """Extract the font scheme from the presentation's theme.
    
    Returns a dict like:
    {
        "major": {"latin": "Calibri", "ea": "宋体", "cs": "..."},
        "minor": {"latin": "Calibri", "ea": "宋体", "cs": "..."}
    }
    """
    theme_fonts = {"major": {}, "minor": {}}
    theme_name = ""
    
    try:
        for slide_master in prs.slide_masters:
            for rel in slide_master.part.rels.values():
                if "theme" in rel.reltype:
                    theme_part = rel.target_part
                    theme_elem = _parse_theme_xml(theme_part)
                    
                    if theme_elem is None:
                        continue
                    
                    font_scheme = theme_elem.find('.//a:fontScheme', NS)
                    if font_scheme is None:
                        results = theme_elem.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}fontScheme')
                        if results:
                            font_scheme = results[0]
                    if font_scheme is None:
                        continue
                    
                    theme_name = font_scheme.get('name', '')
                    
                    for major_minor in ['major', 'minor']:
                        font_group = font_scheme.find(f'a:{major_minor}Font', NS)
                        if font_group is None:
                            font_group = font_scheme.find(qn(f'a:{major_minor}Font'))
                        if font_group is None:
                            continue
                        
                        for font_elem in font_group:
                            tag = font_elem.tag.split('}')[-1] if '}' in font_elem.tag else font_elem.tag
                            typeface = font_elem.get('typeface', '')
                            if typeface:
                                theme_fonts[major_minor][tag] = typeface
                    
                    break
            break
    except Exception as e:
        sys.stderr.write(f"  [WARN] Theme font extraction failed: {e}\n")
    
    if theme_name:
        theme_fonts["_theme_name"] = theme_name
    
    return theme_fonts


def resolve_scheme_color(scheme_color_name, theme_colors):
    """Resolve a scheme color name (like 'dk1', 'accent1') to its hex value."""
    if not theme_colors:
        return None
    return theme_colors.get(scheme_color_name)


def try_get_color_from_fill(fill, theme_colors):
    """Try to extract color from a fill object, handling both RGB and scheme colors."""
    colors = []
    
    # Try direct RGB
    try:
        if fill.type is not None and fill.fore_color:
            try:
                rgb = fill.fore_color.rgb
                if rgb:
                    colors.append(rgb_to_hex(rgb))
            except Exception:
                pass
            
            # Try theme color
            try:
                theme_color = fill.fore_color.theme_color
                if theme_color and theme_colors:
                    # theme_color is an enum like MSO_THEME_COLOR.ACCENT_1
                    # Map enum to theme scheme names
                    tc_map = {
                        0: "dk1",       # DARK_1
                        1: "lt1",       # LIGHT_1
                        2: "dk2",       # DARK_2
                        3: "lt2",       # LIGHT_2
                        4: "accent1",   # ACCENT_1
                        5: "accent2",   # ACCENT_2
                        6: "accent3",   # ACCENT_3
                        7: "accent4",   # ACCENT_4
                        8: "accent5",   # ACCENT_5
                        9: "accent6",   # ACCENT_6
                        10: "hlink",    # HYPERLINK
                        11: "folHlink", # FOLLOWED_HYPERLINK
                    }
                    tc_name = tc_map.get(int(theme_color))
                    if tc_name:
                        resolved = resolve_scheme_color(tc_name, theme_colors)
                        if resolved:
                            colors.append(f"{resolved}/*theme:{tc_name}*/")
            except Exception:
                pass
    except Exception:
        pass
    
    return colors


def try_get_color_from_xml(shape, theme_colors):
    """Low-level XML extraction for colors that python-pptx can't reach."""
    colors = []
    
    try:
        sp = shape._element
        
        # Check solidFill in shape properties
        for fill_elem in sp.iter(qn('a:solidFill')):
            # srgbClr
            srgb = fill_elem.find(qn('a:srgbClr'))
            if srgb is not None:
                val = srgb.get('val')
                if val:
                    colors.append(f"#{val}")
            
            # schemeClr
            scheme = fill_elem.find(qn('a:schemeClr'))
            if scheme is not None:
                val = scheme.get('val')
                if val and theme_colors:
                    resolved = theme_colors.get(val)
                    if resolved:
                        colors.append(f"{resolved}/*theme:{val}*/")
                    else:
                        colors.append(f"/*theme:{val}*/")
    except Exception:
        pass
    
    return colors


def try_get_font_from_run(run, theme_fonts):
    """Get the actual font name from a run, resolving theme fonts."""
    font_name = run.font.name
    font_size = run.font.size
    
    # If font name starts with +, it's a theme font reference
    if font_name and font_name.startswith('+'):
        # +mj-ea → major East Asian, +mn-lt → minor Latin, etc.
        parts = font_name[1:].split('-')  # Remove + and split
        if len(parts) == 2:
            role = parts[0]  # mj or mn
            script = parts[1]  # ea, lt, cs
            
            major_minor = 'major' if role == 'mj' else 'minor'
            script_map = {'ea': 'ea', 'lt': 'latin', 'cs': 'cs'}
            script_key = script_map.get(script, 'latin')
            
            resolved = theme_fonts.get(major_minor, {}).get(script_key)
            if resolved:
                return resolved, font_size, f"{font_name}→{resolved}"
    
    return font_name, font_size, None


# ============================================================
# Core extraction functions
# ============================================================

def extract_colors(slide, theme_colors):
    """Extract color palette from a slide, including theme colors via XML"""
    colors = Counter()

    # Method 1: python-pptx API (works for direct RGB colors)
    # Background color
    bg = slide.background
    if bg and bg.fill:
        for c in try_get_color_from_fill(bg.fill, theme_colors):
            colors[c.replace("/*theme:", " (theme:").replace("*/", ")")] += 5

    # Shape fills and text colors
    for shape in slide.shapes:
        # Skip GroupShapes (no fill attribute)
        shape_type_name = str(shape.shape_type)
        is_group = shape_type_name == 'GROUP (6)' or 'GroupShape' in type(shape).__name__
        
        # Shape fill via python-pptx API
        if not is_group:
            try:
                for c in try_get_color_from_fill(shape.fill, theme_colors):
                    colors[c.replace("/*theme:", " (theme:").replace("*/", ")")] += 3
            except AttributeError:
                pass
            
            # Fallback: XML extraction for shapes where API fails
            for c in try_get_color_from_xml(shape, theme_colors):
                colors[c.replace("/*theme:", " (theme:").replace("*/", ")")] += 3

        # Shape line/stroke
        try:
            if shape.line and shape.line.color:
                try:
                    rgb = shape.line.color.rgb
                    if rgb:
                        colors[rgb_to_hex(rgb)] += 1
                except Exception:
                    pass
                try:
                    tc = shape.line.color.theme_color
                    if tc and theme_colors:
                        tc_map = {0:"dk1",1:"lt1",2:"dk2",3:"lt2",4:"accent1",
                                  5:"accent2",6:"accent3",7:"accent4",8:"accent5",
                                  9:"accent6",10:"hlink",11:"folHlink"}
                        tc_name = tc_map.get(int(tc))
                        if tc_name:
                            resolved = theme_colors.get(tc_name)
                            if resolved:
                                colors[f"{resolved} (theme:{tc_name})"] += 1
                except Exception:
                    pass
        except Exception:
            pass

        # Text colors
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color:
                            try:
                                rgb = run.font.color.rgb
                                if rgb:
                                    colors[rgb_to_hex(rgb)] += 2
                            except Exception:
                                pass
                            try:
                                tc = run.font.color.theme_color
                                if tc and theme_colors:
                                    tc_map = {0:"dk1",1:"lt1",2:"dk2",3:"lt2",4:"accent1",
                                              5:"accent2",6:"accent3",7:"accent4",8:"accent5",
                                              9:"accent6",10:"hlink",11:"folHlink"}
                                    tc_name = tc_map.get(int(tc))
                                    if tc_name:
                                        resolved = theme_colors.get(tc_name)
                                        if resolved:
                                            colors[f"{resolved} (theme:{tc_name})"] += 2
                            except Exception:
                                pass
                    except Exception:
                        pass

    # Method 2: Direct XML scan for schemeClr references (catches what API misses)
    # Map PPTX schemeClr val names to our theme_colors keys
    scheme_alias = {
        'bg1': 'lt1', 'bg2': 'lt2',  # Background
        'tx1': 'dk1', 'tx2': 'dk2',  # Text
    }
    
    try:
        for elem in slide._element.iter():
            tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
            if tag == 'schemeClr':
                val = elem.get('val')
                if val:
                    # Resolve alias first (bg1→lt1, tx1→dk1)
                    resolved_key = scheme_alias.get(val, val)
                    if resolved_key in theme_colors:
                        resolved_hex = theme_colors[resolved_key]
                        colors[f"{resolved_hex} (theme:{val})"] += 2
                    elif val in theme_colors:
                        resolved_hex = theme_colors[val]
                        colors[f"{resolved_hex} (theme:{val})"] += 2
    except Exception:
        pass

    return dict(colors.most_common(30))


def extract_fonts(slide, theme_fonts):
    """Extract font groups from a slide, resolving theme fonts"""
    fonts = {}
    font_sizes = []
    theme_resolution = {}  # Track theme→resolved mappings

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                resolved_name, font_size, resolution = try_get_font_from_run(run, theme_fonts)
                
                # Use resolved name if available, otherwise original
                display_name = resolved_name or run.font.name or "(unnamed)"
                
                if resolution:
                    theme_resolution[run.font.name] = resolution

                if display_name not in fonts:
                    fonts[display_name] = {"count": 0, "sizes": [], "source": "direct"}
                
                fonts[display_name]["count"] += 1
                
                # Mark as theme-derived if resolved
                if resolution:
                    fonts[display_name]["source"] = f"theme({run.font.name})"
                
                if font_size:
                    size_pt = emu_to_pt(font_size)
                    fonts[display_name]["sizes"].append(size_pt)
                    font_sizes.append(size_pt)

    result = {
        "fonts": fonts,
        "size_range": {
            "min": min(font_sizes) if font_sizes else None,
            "max": max(font_sizes) if font_sizes else None
        },
        "theme_resolution": theme_resolution
    }
    return result


def extract_layout_skeletons(slide, slide_idx):
    """Extract layout skeleton (shapes positions/sizes)"""
    skeletons = []

    for shape in slide.shapes:
        skeleton = {
            "shape_id": shape.shape_id,
            "name": shape.name,
            "type": str(shape.shape_type),
            "position": {
                "left_px": emu_to_px(shape.left),
                "top_px": emu_to_px(shape.top),
            },
            "size": {
                "width_px": emu_to_px(shape.width),
                "height_px": emu_to_px(shape.height),
            },
            "has_text": shape.has_text_frame,
        }

        if shape.has_text_frame:
            text = shape.text_frame.text[:100]  # First 100 chars
            skeleton["text_preview"] = text

        skeletons.append(skeleton)

    return skeletons


def extract_decorations(slide, theme_colors):
    """Extract decorative elements (lines, logos, corner marks)"""
    decorations = []

    for shape in slide.shapes:
        # Skip GroupShapes
        is_group = 'GroupShape' in type(shape).__name__
        
        # Skip shapes with significant text (those are content, not decoration)
        if shape.has_text_frame and len(shape.text_frame.text.strip()) > 20:
            continue

        # Small shapes are likely decorative
        width_px = emu_to_px(shape.width) or 0
        height_px = emu_to_px(shape.height) or 0

        is_small = width_px < 100 or height_px < 100
        is_line = width_px > 500 and height_px < 5
        is_vline = height_px > 500 and width_px < 5

        if is_small or is_line or is_vline:
            decoration = {
                "type": "line" if (is_line or is_vline) else "shape",
                "name": shape.name,
                "position": {
                    "left_px": emu_to_px(shape.left),
                    "top_px": emu_to_px(shape.top),
                },
                "size": {
                    "width_px": width_px,
                    "height_px": height_px,
                }
            }

            # Try to get fill color (API first, then XML) — skip for GroupShapes
            if not is_group:
                try:
                    fill_colors = try_get_color_from_fill(shape.fill, theme_colors)
                    if fill_colors:
                        decoration["fill_color"] = fill_colors[0]
                except AttributeError:
                    pass
                
                xml_colors = try_get_color_from_xml(shape, theme_colors)
                if xml_colors and "fill_color" not in decoration:
                    decoration["fill_color"] = xml_colors[0]

            decorations.append(decoration)

    return decorations


def extract_placeholders(slide):
    """Extract image/text placeholder rules"""
    placeholders = []

    for shape in slide.shapes:
        # Image placeholders
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            placeholders.append({
                "type": "image",
                "name": shape.name,
                "aspect_ratio": round(
                    (shape.width / shape.height) if shape.height else 0, 2
                ),
                "width_px": emu_to_px(shape.width),
                "height_px": emu_to_px(shape.height),
            })

        # Text placeholders (large text boxes are likely content placeholders)
        elif shape.has_text_frame:
            text_len = len(shape.text_frame.text.strip())
            width_px = emu_to_px(shape.width) or 0
            if width_px > 400 and text_len > 10:
                placeholders.append({
                    "type": "text",
                    "name": shape.name,
                    "text_length": text_len,
                    "width_px": width_px,
                    "height_px": emu_to_px(shape.height),
                    "capacity_estimate": round(width_px * (emu_to_px(shape.height) or 0) / 1000, 1)
                })

    return placeholders


def classify_rhythm(density, slide_idx, total_slides):
    """Classify page rhythm based on content density"""
    if slide_idx == 0:
        return "anchor"  # Cover
    if slide_idx == total_slides - 1:
        return "anchor"  # End
    if density < 0.3:
        return "breathing"  # Low density = visual impact page
    return "dense"  # Default: information dense


def compute_density(slide, canvas_area):
    """Compute content density of a slide"""
    text_area = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Rough text area estimate
            w = emu_to_px(shape.width) or 0
            h = emu_to_px(shape.height) or 0
            text_area += w * h

    if canvas_area == 0:
        return 0
    return min(text_area / canvas_area, 1.0)


def dissect_pptx(pptx_path, project_path):
    """Main dissection function"""
    pptx_path = Path(pptx_path)
    project_path = Path(project_path)

    if not pptx_path.exists():
        sys.stderr.write(f"ERROR: PPTX file not found: {pptx_path}\n")
        sys.exit(1)

    # Ensure output directories
    elements_dir = project_path / "template_elements"
    elements_dir.mkdir(parents=True, exist_ok=True)

    # Load presentation
    prs = Presentation(str(pptx_path))

    # Canvas dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    canvas_w_px = emu_to_px(slide_width)
    canvas_h_px = emu_to_px(slide_height)
    canvas_area = canvas_w_px * canvas_h_px

    total_slides = len(prs.slides)

    # ============================================================
    # NEW: Extract theme data first (shared across all slides)
    # ============================================================
    sys.stderr.write("  Extracting theme colors and fonts...\n")
    theme_colors = extract_theme_colors(prs)
    theme_fonts = extract_theme_fonts(prs)
    
    sys.stderr.write(f"  Theme colors: {len(theme_colors)} entries\n")
    for k, v in list(theme_colors.items())[:5]:
        sys.stderr.write(f"    {k}: {v}\n")
    
    sys.stderr.write(f"  Theme fonts: major={theme_fonts.get('major', {})}, minor={theme_fonts.get('minor', {})}\n")

    # Save theme data
    theme_data = {
        "source": str(pptx_path.name),
        "colors": theme_colors,
        "fonts": theme_fonts,
    }
    theme_path = elements_dir / "theme.json"
    with open(theme_path, "w", encoding="utf-8") as f:
        json.dump(theme_data, f, ensure_ascii=False, indent=2)
    sys.stderr.write(f"  Written: theme.json\n")

    # ============================================================
    # Per-slide extraction (with theme resolution)
    # ============================================================
    all_colors = Counter()
    all_fonts = {}
    layout_skeletons = {}
    all_decorations = []
    all_placeholders = []
    rhythm_map = {}
    all_theme_resolution = {}

    for slide_idx, slide in enumerate(prs.slides):
        page_id = f"P{slide_idx + 1:02d}"

        # 1. Colors (with theme resolution)
        page_colors = extract_colors(slide, theme_colors)
        all_colors.update(page_colors)

        # 2. Fonts (with theme resolution)
        page_fonts = extract_fonts(slide, theme_fonts)
        for font_name, font_data in page_fonts.get("fonts", {}).items():
            if font_name not in all_fonts:
                all_fonts[font_name] = {"count": 0, "sizes": [], "source": font_data.get("source", "direct")}
            all_fonts[font_name]["count"] += font_data["count"]
            all_fonts[font_name]["sizes"].extend(font_data.get("sizes", []))
        
        for k, v in page_fonts.get("theme_resolution", {}).items():
            all_theme_resolution[k] = v

        # 3. Layout skeletons
        layout_skeletons[page_id] = extract_layout_skeletons(slide, slide_idx)

        # 4. Decorations (with theme resolution)
        page_decs = extract_decorations(slide, theme_colors)
        all_decorations.extend(page_decs)

        # 5. Placeholders
        page_phs = extract_placeholders(slide)
        all_placeholders.extend(page_phs)

        # 6. Rhythm
        density = compute_density(slide, canvas_area)
        rhythm_map[page_id] = {
            "density": round(density, 3),
            "rhythm": classify_rhythm(density, slide_idx, total_slides)
        }

    # --- Build color palette ---
    top_colors = all_colors.most_common(20)
    color_palette = {
        "source": str(pptx_path.name),
        "canvas": {"width_px": canvas_w_px, "height_px": canvas_h_px},
        "extracted_colors": [
            {"hex": c, "weight": w} for c, w in top_colors
        ],
        "theme_colors": theme_colors,
        "recommendation": {}
    }

    # Identify 60-30-10 roles — use all colors sorted by weight (already sorted)
    # Clean up hex values (remove theme annotation for recommendation)
    def clean_hex(h):
        return h.split(" (theme:")[0].strip() if " (theme:" in h else h
    
    # Filter out empty/broken entries
    valid_colors = [(c, w) for c, w in top_colors if clean_hex(c)]
    
    # Identify background (usually white/light) vs primary (usually dark accent)
    # Strategy: the highest-weight color is often background, second is primary
    # But if accent1 (浙大蓝) dominates, it should be primary
    if len(valid_colors) >= 2:
        bg_candidates = [c for c, w in valid_colors if clean_hex(c).upper() in ('#FFFFFF', '#E7E6E6', '#F2F2F2')]
        primary_candidates = [c for c, w in valid_colors if "(theme:accent" in c and clean_hex(c) not in ('#FFFFFF',)]
        
        if bg_candidates:
            color_palette["recommendation"]["background"] = clean_hex(bg_candidates[0])
        else:
            color_palette["recommendation"]["background"] = clean_hex(valid_colors[0][0])
        
        if primary_candidates:
            # Sort by weight
            primary_sorted = sorted(
                [(c, dict(top_colors).get(c, 0)) for c in primary_candidates],
                key=lambda x: x[1], reverse=True
            )
            color_palette["recommendation"]["primary"] = clean_hex(primary_sorted[0][0])
            if len(primary_sorted) > 1:
                color_palette["recommendation"]["accent"] = clean_hex(primary_sorted[1][0])
            else:
                color_palette["recommendation"]["accent"] = "#1A3A5C"
        else:
            color_palette["recommendation"]["primary"] = clean_hex(valid_colors[1][0])
            color_palette["recommendation"]["accent"] = clean_hex(valid_colors[2][0]) if len(valid_colors) > 2 else "#1A3A5C"
        
        color_palette["recommendation"]["body_text"] = "#1D1D1F"
        color_palette["recommendation"]["secondary_text"] = "#545458"

    # --- Build font groups ---
    font_groups = {
        "source": str(pptx_path.name),
        "theme_fonts": theme_fonts,
        "theme_resolution": all_theme_resolution,
        "fonts": {}
    }

    # Sort by usage count
    sorted_fonts = sorted(all_fonts.items(), key=lambda x: x[1]["count"], reverse=True)
    for font_name, font_data in sorted_fonts[:15]:
        sizes = font_data.get("sizes", [])
        font_groups["fonts"][font_name] = {
            "count": font_data["count"],
            "size_range": f"{min(sizes):.0f}-{max(sizes):.0f}pt" if sizes else "unknown",
            "recommended_role": "title" if font_data["count"] <= 5 else ("body" if font_data["count"] <= 20 else "secondary"),
            "source": font_data.get("source", "direct")
        }

    # --- Write 7 JSON files (6 + theme.json already written) ---
    outputs = {
        "color_palette.json": color_palette,
        "font_groups.json": font_groups,
        "layout_skeletons.json": layout_skeletons,
        "decorations.json": {"source": str(pptx_path.name), "items": all_decorations},
        "placeholders.json": {"source": str(pptx_path.name), "items": all_placeholders},
        "rhythm_map.json": rhythm_map,
    }

    for filename, data in outputs.items():
        outpath = elements_dir / filename
        with open(outpath, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        sys.stderr.write(f"  Written: {outpath.name}\n")

    # --- Update spec_lock.md ---
    spec_lock_path = project_path / "spec_lock.md"
    if spec_lock_path.exists():
        with open(spec_lock_path, "r", encoding="utf-8") as f:
            content = f.read()
    else:
        content = "# Spec Lock — Machine-readable execution contract\n\n"

    # Append template dissection section
    dissect_section = f"""
## Template Dissection Results

> Source: {pptx_path.name}
> Canvas: {canvas_w_px}x{canvas_h_px}px

### Theme Colors
{json.dumps(theme_colors, ensure_ascii=False, indent=2)}

### Theme Fonts
{json.dumps(theme_fonts, ensure_ascii=False, indent=2)}

### Color Palette (Recommendation)
{json.dumps(color_palette.get("recommendation", {}), ensure_ascii=False, indent=2)}

### Font Groups
{json.dumps({k: v for k, v in list(font_groups["fonts"].items())[:8]}, ensure_ascii=False, indent=2)}

### Rhythm Map
{json.dumps(rhythm_map, ensure_ascii=False, indent=2)}
"""

    with open(spec_lock_path, "w", encoding="utf-8") as f:
        f.write(content + dissect_section)
    sys.stderr.write(f"  Updated: spec_lock.md\n")

    # --- Update design_spec.md ---
    design_spec_path = project_path / "design_spec.md"
    if design_spec_path.exists():
        with open(design_spec_path, "r", encoding="utf-8") as f:
            ds_content = f.read()
    else:
        ds_content = f"# {project_path.name} - Design Spec\n\n"

    # Append human-readable template info
    template_section = f"""
## Template Analysis

> Source: {pptx_path.name}
> Canvas: {canvas_w_px}x{canvas_h_px}px ({canvas_w_px/96:.1f}x{canvas_h_px/96:.1f} inches)

### Theme Colors
| Scheme Name | HEX |
|-------------|-----|
"""
    for name, hex_val in theme_colors.items():
        template_section += f"| {name} | `{hex_val}` |\n"

    template_section += f"""
### Theme Fonts
| Role | Latin | East Asian | CS |
|------|-------|------------|----|
| Major (Heading) | {theme_fonts.get('major',{}).get('latin','-')} | {theme_fonts.get('major',{}).get('ea','-')} | {theme_fonts.get('major',{}).get('cs','-')} |
| Minor (Body) | {theme_fonts.get('minor',{}).get('latin','-')} | {theme_fonts.get('minor',{}).get('ea','-')} | {theme_fonts.get('minor',{}).get('cs','-')} |

### Extracted Color Palette
| Role | HEX | Source |
|------|-----|--------|
"""
    for role, hex_val in color_palette.get("recommendation", {}).items():
        template_section += f"| {role} | `{hex_val}` | Template extraction |\n"

    template_section += f"""
### Extracted Fonts
| Font | Usage Count | Size Range | Source | Recommended Role |
|------|------------|------------|--------|-----------------|
"""
    for font_name, font_data in sorted_fonts[:10]:
        template_section += f"| {font_name} | {font_data['count']} | {font_data.get('size_range', '?')} | {font_data.get('source', '?')} | {font_groups['fonts'].get(font_name, {}).get('recommended_role', 'unknown')} |\n"

    template_section += f"""
### Rhythm Map
| Page | Density | Rhythm |
|------|---------|--------|
"""
    for page_id, rhythm_data in rhythm_map.items():
        template_section += f"| {page_id} | {rhythm_data['density']:.3f} | {rhythm_data['rhythm']} |\n"

    with open(design_spec_path, "w", encoding="utf-8") as f:
        f.write(ds_content + template_section)
    sys.stderr.write(f"  Updated: design_spec.md\n")

    # --- Summary ---
    sys.stderr.write(f"\n=== Template Dissection Complete (v2) ===\n")
    sys.stderr.write(f"  Source: {pptx_path.name}\n")
    sys.stderr.write(f"  Slides: {total_slides}\n")
    sys.stderr.write(f"  Canvas: {canvas_w_px}x{canvas_h_px}px\n")
    sys.stderr.write(f"  Theme colors: {len(theme_colors)}\n")
    sys.stderr.write(f"  Theme fonts: major={theme_fonts.get('major',{})}, minor={theme_fonts.get('minor',{})}\n")
    sys.stderr.write(f"  Unique extracted colors: {len(all_colors)}\n")
    sys.stderr.write(f"  Unique fonts (resolved): {len(all_fonts)}\n")
    sys.stderr.write(f"  Theme resolution map: {all_theme_resolution}\n")
    sys.stderr.write(f"  Decorations: {len(all_decorations)}\n")
    sys.stderr.write(f"  Rhythm: {sum(1 for v in rhythm_map.values() if v['rhythm']=='anchor')} anchor, {sum(1 for v in rhythm_map.values() if v['rhythm']=='dense')} dense, {sum(1 for v in rhythm_map.values() if v['rhythm']=='breathing')} breathing\n")
    sys.stderr.write(f"  Output: {elements_dir}/\n")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        sys.stderr.write("Usage: python template_dissector.py <pptx_path> <project_path>\n")
        sys.exit(1)

    dissect_pptx(sys.argv[1], sys.argv[2])
